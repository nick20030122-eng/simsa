import ast
import io
import re
from typing import Any, Dict, List


def extract_pdf_text(uploaded_file) -> str:
    """
    Streamlit UploadedFile(.pdf)에서 페이지별 텍스트를 추출합니다.
    `pypdf` 패키지가 필요합니다 (`pip install pypdf`).
    """
    if uploaded_file is None:
        return ""
    raw = uploaded_file.getvalue()
    if not raw:
        return ""

    try:
        from pypdf import PdfReader
    except ImportError as e:
        raise ImportError("PDF 처리를 위해 pypdf 패키지가 필요합니다. pip install pypdf") from e

    reader = PdfReader(io.BytesIO(raw))
    parts: List[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            t = (page.extract_text() or "").strip()
        except Exception:
            t = ""
        if t:
            parts.append(f"[PDF 페이지 {i}]\n{t}")
    return "\n\n".join(parts).strip()


def extract_ppt_text(uploaded_file) -> str:
    """
    Streamlit UploadedFile(.pptx)에서 텍스트를 추출합니다.
    - python-pptx가 없거나 실패하면 예외를 그대로 발생시킵니다.
    """
    if uploaded_file is None:
        return ""
    raw = uploaded_file.getvalue()
    if not raw:
        return ""

    # python-pptx는 설치돼 있어야 합니다.
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    lines: List[str] = []
    for s_i, slide in enumerate(prs.slides, start=1):
        lines.append(f"[슬라이드 {s_i}]")
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                lines.append(text)
    return "\n".join(lines).strip()


def extract_ppt_signals(ppt_text: str) -> Dict[str, Any]:
    """
    '파급효과' 점수에 반영하기 위해 PPT에서 우선적으로 찾아낼 신호를 추출합니다.
    - 문제 정의(problem/pain)
    - 시장 규모(market size/TAM/SAM/SOM/성장률/매출 등)
    """
    text = ppt_text or ""
    text = re.sub(r"\s+", " ", text)

    problem_keywords = [
        r"문제\s*정의",
        r"문제점",
        r"페인\s*포인트",
        r"pain\s*point",
        r"문제",
        r"고객\s*불편",
        r"니즈",
        r"needs?",
    ]
    market_keywords = [
        r"시장\s*규모",
        r"TAM",
        r"SAM",
        r"SOM",
        r"Market\s*Size",
        r"시장",
        r"CAGR",
        r"성장률",
        r"매출",
        r"거래액",
        r"MAU",
        r"사용자\s*수",
        r"억원|만\s*원|달러|\$|USD|KRW",
        r"%|퍼센트",
    ]

    def _find_snippets(patterns: List[str], max_snippets: int) -> List[str]:
        out: List[str] = []
        for p in patterns:
            for m in re.finditer(p, text, flags=re.IGNORECASE):
                start = max(0, m.start() - 60)
                end = min(len(text), m.end() + 120)
                snippet = text[start:end].strip()
                if snippet and snippet not in out:
                    out.append(snippet)
                if len(out) >= max_snippets:
                    return out
        return out

    problem = _find_snippets(problem_keywords, max_snippets=5)
    market = _find_snippets(market_keywords, max_snippets=7)

    return {
        "problem_definition_snippets": problem,
        "market_size_snippets": market,
    }


def _decode_bytes(raw: bytes) -> str:
    if not raw:
        return ""
    for enc in ("utf-8", "cp949", "utf-16"):
        try:
            return raw.decode(enc)
        except Exception:
            continue
    return raw.decode("utf-8", errors="replace")


def _python_complexity_metrics(tree: ast.AST) -> Dict[str, int]:
    counts = {
        "functions": 0,
        "classes": 0,
        "imports": 0,
        "branches": 0,
        "loops": 0,
        "exceptions": 0,
        "bool_ops": 0,
    }
    for node in ast.walk(tree):
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
            counts["functions"] += 1
        elif isinstance(node, ast.ClassDef):
            counts["classes"] += 1
        elif isinstance(node, (ast.Import, ast.ImportFrom)):
            counts["imports"] += 1
        elif isinstance(node, (ast.If, ast.Match)):
            counts["branches"] += 1
        elif isinstance(node, (ast.For, ast.AsyncFor, ast.While)):
            counts["loops"] += 1
        elif isinstance(node, (ast.Try, ast.Raise)):
            counts["exceptions"] += 1
        elif isinstance(node, ast.BoolOp):
            counts["bool_ops"] += 1
    return counts


def _python_imports(tree: ast.AST) -> List[str]:
    libs: List[str] = []
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for n in node.names:
                root = (n.name or "").split(".")[0]
                if root and root not in libs:
                    libs.append(root)
        elif isinstance(node, ast.ImportFrom):
            mod = (node.module or "").split(".")[0]
            if mod and mod not in libs:
                libs.append(mod)
    return libs


def _js_imports(text: str) -> List[str]:
    libs: List[str] = []
    patterns = [
        r"""import\s+.*?\s+from\s+["']([^"']+)["']""",
        r"""require\(\s*["']([^"']+)["']\s*\)""",
    ]
    for pat in patterns:
        for m in re.finditer(pat, text):
            name = (m.group(1) or "").strip()
            root = name.split("/")[0] if name else ""
            if root and root not in libs and not root.startswith("."):
                libs.append(root)
    return libs


def _js_complexity_heuristic(text: str) -> Dict[str, int]:
    return {
        "functions": len(re.findall(r"\bfunction\b|\=\>\s*\{", text)),
        "classes": len(re.findall(r"\bclass\b", text)),
        "imports": len(re.findall(r"\bimport\b|\brequire\b", text)),
        "branches": len(re.findall(r"\bif\b|\bswitch\b|\bcase\b", text)),
        "loops": len(re.findall(r"\bfor\b|\bwhile\b|\bforEach\b", text)),
        "exceptions": len(re.findall(r"\btry\b|\bcatch\b|\bthrow\b", text)),
    }


def analyze_code_files(uploaded_files) -> Dict[str, Any]:
    """
    코드 파일에서 '개발 난이도' 점수에 도움되는 정보를 요약합니다.
    - 사용 라이브러리(Imports)
    - 대략적인 복잡도(구조/분기/루프/예외 처리 등)
    """
    files_out: List[Dict[str, Any]] = []
    all_libs: List[str] = []

    for f in uploaded_files or []:
        name = (f.name or "file").strip()
        raw = f.getvalue() or b""
        text = _decode_bytes(raw)
        ext = name.lower().split(".")[-1] if "." in name else ""

        info: Dict[str, Any] = {"name": name, "ext": ext, "size_bytes": len(raw)}

        if ext == "py":
            try:
                tree = ast.parse(text)
                info["imports"] = _python_imports(tree)
                info["metrics"] = _python_complexity_metrics(tree)
            except Exception as e:
                info["parse_error"] = str(e)
                info["imports"] = []
                info["metrics"] = {}
        elif ext in {"js", "jsx", "ts", "tsx"}:
            info["imports"] = _js_imports(text)
            info["metrics"] = _js_complexity_heuristic(text)
        else:
            info["imports"] = _js_imports(text)
            info["metrics"] = {"imports": len(info["imports"])}

        for lib in info.get("imports", []) or []:
            if lib not in all_libs:
                all_libs.append(lib)

        files_out.append(info)

    return {
        "libraries_detected": all_libs,
        "files": files_out,
    }

